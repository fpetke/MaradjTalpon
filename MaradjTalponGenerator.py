import pptx
from os import listdir
from typing import List, Dict, Tuple
import numpy as np
import sys
import string
import math

class MaradjTalpon:
    question_dir: str = "kerdesek"
    questions: List[str] = None
    answers: List[str] = None
    categories: Dict[str, Tuple[int]] = None
    question_ids: np.ndarray = None


    def __init__(self, max_num_questions: int, question_dir = None) -> None:
        if not question_dir is None:
            self.question_dir = question_dir

        self.load_questions_answers()
        self.max_num_questions = (max_num_questions if max_num_questions < self.question_ids.size else self.question_ids.size)

        self.root = pptx.Presentation()
        self.create_title_slide()
        self.generate_question_slides()

    def load_questions_answers(self) -> None:
        """
        This function loads questions and answers from files in a directory, assigns categories to them,
        and shuffles the question IDs.
        """
        self.questions = []
        self.answers = []
        self.categories = {}

        for fname in listdir(self.question_dir):
            category_start = len(self.answers)

            with open(self.question_dir+'/'+fname, 'r', encoding="utf-8") as f:
                append_to_questions = True
                for line in f:
                    if append_to_questions:
                        self.questions.append(line.strip())
                        append_to_questions = False
                    else:
                        self.answers.append(line.strip())
                        append_to_questions = True

            self.categories[fname] = (category_start, len(self.answers)-1)
                        
        self.question_ids = np.arange(0, len(self.questions), dtype=int)
        np.random.shuffle(self.question_ids)

    def get_category_of_question(self, id: int) -> str | None:
        """
        This function takes an ID as input and returns the category of the question based on the ID
        range specified in the categories dictionary.
        
        :param id: The `id` parameter in the `get_category_of_question` method is an integer
        representing the unique identifier of a question. The method iterates through the categories
        stored in the object instance and determines which category the question with the given `id`
        falls into based on the range specified for each category
        :type id: int
        :return: The function `get_category_of_question` returns a string representing the category of a
        question based on its ID. If the ID falls within the range specified for a category in the
        `self.categories` dictionary, the function returns the category. If no category is found for the
        given ID, it returns `None`.
        """
        for category in self.categories:
            if id >= self.categories[category][0] and id <= self.categories[category][1]:
                return category
        return None

    def generate_question_slides(self) -> None:
        """
        The function generates question and answer slides for a specified number of questions.
        """
        for question_id in self.question_ids[:self.max_num_questions]:
            self.create_question_slide(question_id)
            self.create_answer_slide(question_id)

    def create_title_slide(self) -> None:
        """
        The function `create_title_slide` creates a title slide with the text "Maradj Talpon!" and the
        number of questions specified in the `max_num_questions` variable.
        """
        slide = self.root.slides.add_slide(self.root.slide_layouts[0])
        slide.shapes.title.text = "Maradj Talpon!"
        slide.placeholders[1].text = "%d izgalmas kérdés különféle kategóriákból." %(self.max_num_questions)

    def create_question_slide(self, question_id: int) -> None:
        """
        The function creates a question slide by replacing a portion of the text with underscores and
        underlines in a PowerPoint presentation.
        
        :param question_id: The `question_id` parameter is an integer that represents the unique
        identifier of the question for which you want to create a question slide. This identifier is
        used to retrieve the question content and format it into a slide with certain text runs modified
        to be blank for the user to fill in
        :type question_id: int
        """
        txFrame = self.create_answer_slide(question_id)

        runs_to_choose = []
        for run in txFrame.paragraphs[0].runs:
            if not (run.text in string.punctuation or run.text == '⎵' or run.text == ' '):
                runs_to_choose.append(run)

        ratio_to_make_blank = 0.6
        N_to_choose = np.max([1, int(math.ceil(len(runs_to_choose) * ratio_to_make_blank))])
        if N_to_choose > len(runs_to_choose):
            N_to_choose = len(runs_to_choose)
        chosen_runs = np.random.choice(runs_to_choose, N_to_choose, replace=False)

        for run in chosen_runs:
            run.text = '_'
            run.font.underline = pptx.enum.text.MSO_UNDERLINE.SINGLE_LINE
            run.font.color.rgb = pptx.dml.color.RGBColor(242, 7, 7)

    def create_answer_slide(self, question_id: int) -> pptx.text.text.TextFrame:
        """
        This function creates a slide in a PowerPoint presentation with a question and its corresponding
        answers formatted with specific styles.
        
        :param question_id: The `question_id` parameter in the `create_answer_slide` method is used to
        specify which question's answer slide should be created. This parameter is an integer that
        corresponds to the index of the question in the `questions` and `answers` lists within the class
        instance. By providing the `question
        :type question_id: int
        :return: The `create_answer_slide` method returns a `pptx.text.text.TextFrame` object that
        contains the answer text for a given question ID.
        """
        slide = self.root.slides.add_slide(self.root.slide_layouts[5])
        slide.shapes.title.text = self.questions[question_id]

        self.add_category_to_slide(slide, question_id)
        self.add_slide_numbering(slide, question_id)

        width = self.root.slide_width * 0.8
        height = self.root.slide_height * 0.4
        left = (self.root.slide_width - width) / 2
        top = self.root.slide_height * 0.8 - height

        txFrame = slide.shapes.add_textbox(left, top, width, height).text_frame
        txFrame.word_wrap = True

        p = txFrame.paragraphs[0]
        p.alignment = pptx.enum.text.PP_ALIGN.CENTER
        p.font.size = pptx.util.Pt(40)
        p.font.name = "Consolas"

        for i in range(len(self.answers[question_id])):
            if i > 0:
                run = p.add_run()
                run.text = " "

            run = p.add_run()
            run.text = self.answers[question_id][i]

            if run.text == ' ':
                run.text = '⎵'
                run.font.color.rgb = pptx.dml.color.RGBColor(199, 191, 191)
            elif not run.text in string.punctuation:
                run.font.underline = pptx.enum.text.MSO_UNDERLINE.SINGLE_LINE

        return txFrame

    def add_category_to_slide(self, slide, question_id: int) -> None:
        """
        This function adds a category label to a slide in a PowerPoint presentation based on a given
        question ID.
        
        :param slide: The `slide` parameter in the `add_category_to_slide` method represents the slide
        object to which you want to add a category label for a specific question. This method adds a
        text box to the slide with the category information of the question specified by the
        `question_id`
        :param question_id: The `question_id` parameter is an integer that represents the unique
        identifier of a question. It is used to retrieve the category of the question and add it to a
        specific slide in a PowerPoint presentation
        :type question_id: int
        """
        category = self.get_category_of_question(question_id)

        width = pptx.util.Inches(3)
        height = pptx.util.Pt(18)
        left = self.root.slide_width * 0.025
        top = self.root.slide_height * 0.975 - height

        txFrame = slide.shapes.add_textbox(left, top, width, height).text_frame
        txFrame.text = category + " - %d" %(question_id - self.categories[category][0] + 1)
        txFrame.fit_text()

    def add_slide_numbering(self, slide, question_id: int):
        """
        The function `add_slide_numbering` adds slide numbering to a PowerPoint slide based on a given
        question ID.
        
        :param slide: The `slide` parameter is the slide object to which you want to add the slide
        numbering
        :param question_id: The `question_id` parameter is an integer that represents the unique
        identifier of a question. It is used to identify the specific question for which the slide
        numbering is being added in the PowerPoint presentation
        :type question_id: int
        """
        width = pptx.util.Inches(3)
        height = pptx.util.Pt(18)
        left = self.root.slide_width * 0.975 - width
        top = self.root.slide_height * 0.975 - height

        txFrame = slide.shapes.add_textbox(left, top, width, height).text_frame
        p = txFrame.paragraphs[0]
        p.alignment = pptx.enum.text.PP_ALIGN.RIGHT
        p.text = "%d / %d" %(np.where(self.question_ids==question_id)[0][0]+1, self.max_num_questions)
        txFrame.fit_text()

    def save_pptx(self, fname) -> None:
        """
        The function `save_pptx` saves a PowerPoint file using the specified filename.
        
        :param fname: The `fname` parameter in the `save_pptx` method is a string that represents the
        file name or path where the PowerPoint presentation (PPTX file) will be saved
        """
        self.root.save(fname)


def main():
    max_num_questions = 1000 # just in case it wasn't specified as a command line argument
    pres = MaradjTalpon(int(sys.argv[1]) if len(sys.argv) == 2 else max_num_questions)

    pres.save_pptx('MaradjTalpon.pptx')

if __name__ == "__main__":
    main()
